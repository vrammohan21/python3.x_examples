from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Theme colors and font
primary_color = RGBColor(30, 64, 175)  # Indigo
accent_color = RGBColor(16, 185, 129)  # Emerald
font_name = "Segoe UI"

# Utility Functions
def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    # Style
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = accent_color
    return slide

def add_bullet_slide(title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text
    content.text = "\n".join([f"• {b}" for b in bullets])
    # Style
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    return slide

# ------------------------
# Build the 20 Slides
# ------------------------

# 1. Title Slide
add_title_slide("Java & Python Full-Stack Development", "Course Launch Showcase\nEmpowering Future Developers")

# 2. Agenda
add_bullet_slide("Agenda", [
    "Welcome & Introduction",
    "Why Full-Stack Development?",
    "Program Structure & Tracks",
    "Curriculum Overview",
    "Projects & Capstone",
    "Career Pathways",
    "Q&A"
])

# 3. Why Full-Stack Development
add_bullet_slide("Why Full-Stack Development?", [
    "High demand for end-to-end developers",
    "Work on frontend + backend",
    "Better salaries & career growth",
    "Complete ownership of projects"
])

# 4. Program Overview
add_bullet_slide("Program Overview", [
    "Two tracks: Java FSD & Python FSD",
    "Frontend, Backend, Database, DevOps",
    "Hands-on projects every module",
    "Capstone project with real-world scenario"
])

# 5. Java Full-Stack Track
add_bullet_slide("Java Full-Stack Track", [
    "Core Java, Collections, Streams",
    "Spring Boot, REST APIs, JPA",
    "Frontend: HTML, CSS, JS, React/Angular",
    "Database: MySQL/PostgreSQL, MongoDB"
])

# 6. Python Full-Stack Track
add_bullet_slide("Python Full-Stack Track", [
    "Python, OOP, venv",
    "Django/Flask, REST APIs",
    "Frontend: React",
    "Database: PostgreSQL, MongoDB"
])

# 7. Frontend Foundations
add_bullet_slide("Frontend Foundations", [
    "HTML5, CSS3, JavaScript ES6+",
    "Responsive design with Flexbox & Grid",
    "React or Angular for dynamic UIs",
    "State management basics"
])

# 8. Java Backend Deep Dive
add_bullet_slide("Java Backend Deep Dive", [
    "Spring Boot fundamentals",
    "RESTful services & validation",
    "Data handling with JPA/Hibernate",
    "Security: JWT, OAuth2"
])

# 9. Python Backend Deep Dive
add_bullet_slide("Python Backend Deep Dive", [
    "Django MVC structure",
    "Django REST Framework",
    "Flask microservices",
    "Testing & coverage"
])

# 10. Databases
add_bullet_slide("Databases & Persistence", [
    "Relational: MySQL/PostgreSQL",
    "ORM usage",
    "NoSQL with MongoDB",
    "Migrations & data seeding"
])

# 11. DevOps & Deployment
add_bullet_slide("DevOps & Deployment", [
    "Git, GitHub Actions, CI/CD",
    "Docker for containerization",
    "Kubernetes basics",
    "Monitoring & logs"
])

# 12. Tools & IDEs
add_bullet_slide("Tools & IDEs", [
    "IntelliJ IDEA, PyCharm, VS Code",
    "Linters & formatters",
    "Package managers (npm, pip)",
    "Debugging & profiling"
])

# 13. Projects by Module
add_bullet_slide("Projects by Module", [
    "Static site (HTML/CSS)",
    "React SPA with routing",
    "Java Spring Boot REST service",
    "Python Django REST API"
])

# 14. Capstone Project
add_bullet_slide("Capstone Project", [
    "Real-world full-stack app",
    "Auth, CRUD, API integration",
    "CI/CD pipeline",
    "Deployed to cloud"
])

# 15. Schedule Overview
add_bullet_slide("Program Schedule (16 Weeks)", [
    "W1-2: Frontend",
    "W3-4: React",
    "W5-8: Backend (Java/Python)",
    "W9-10: Databases",
    "W11-12: DevOps",
    "W13-16: Capstone"
])

# 16. Assessments
add_bullet_slide("Assessments & Certification", [
    "Quizzes & coding challenges",
    "Project evaluations",
    "Capstone review",
    "Completion certificate"
])

# 17. Career Pathways
add_bullet_slide("Career Pathways", [
    "Full-Stack Developer",
    "Backend Engineer",
    "Frontend Engineer",
    "DevOps Engineer"
])

# 18. Portfolio Building
add_bullet_slide("Portfolio Building", [
    "GitHub repositories",
    "Professional README",
    "CI badges",
    "Project documentation"
])

# 19. Support & Mentorship
add_bullet_slide("Support & Mentorship", [
    "Live sessions & recordings",
    "1:1 doubt clearing",
    "Community discussions",
    "Career guidance"
])

# 20. Call to Action
""" add_bullet_slide("Join the Next Cohort", [
    "Apply now – Limited seats",
    "Email: admissions@example.com",
    "Phone: +91-00000 00000",
    "Website: example.com/fsd"
]) """

# Save the PPT file
prs.save("Java_Python_FSD_Course_Detailed.pptx")
print("PPT Generated: Java_Python_FSD_Course_Detailed.pptx")
